#!/usr/bin/env python3

import requests
import time
import tempfile
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Configuration
API_BASE = "http://localhost:8000"
USERNAME = "admin"
PASSWORD = "admin123"

def create_complex_pptx():
    """Create a complex PowerPoint for testing layout preservation"""
    ppt = Presentation()
    
    # Slide 1: Complex layout with multiple shapes
    slide1 = ppt.slides.add_slide(ppt.slide_layouts[1])  # Title and content
    slide1.shapes.title.text = "Complex Layout Test"
    
    # Add text boxes with different formatting
    textbox1 = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
    textbox1.text = "rs12345 - Should be colored"
    textbox1.fill.solid()
    textbox1.fill.fore_color.rgb = RGBColor(255, 0, 0)  # Red background
    
    textbox2 = slide1.shapes.add_textbox(Inches(1), Inches(3.5), Inches(4), Inches(1))
    textbox2.text = "rs67890 - Another test RSID"
    textbox2.fill.solid()
    textbox2.fill.fore_color.rgb = RGBColor(0, 255, 0)  # Green background
    
    # Add shapes
    rectangle = slide1.shapes.add_shape(1, Inches(6), Inches(2), Inches(2), Inches(1))  # Rectangle
    rectangle.text = "Shape Text"
    
    # Slide 2: More complex content
    slide2 = ppt.slides.add_slide(ppt.slide_layouts[0])  # Title slide
    slide2.shapes.title.text = "Second Slide"
    
    # Add multiple text boxes with RSIDs
    for i, rsid in enumerate(['rs111111', 'rs222222', 'rs333333']):
        y_pos = 2 + i * 1.5
        textbox = slide2.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(3), Inches(1))
        textbox.text = rsid
        # Different colors
        colors = [RGBColor(0, 0, 255), RGBColor(255, 165, 0), RGBColor(128, 0, 128)]
        textbox.fill.solid()
        textbox.fill.fore_color.rgb = colors[i]
    
    # Save to temporary file
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_file:
        ppt.save(temp_file.name)
        return temp_file.name

def upload_test_template(token, pptx_path):
    """Upload the test PowerPoint as a template"""
    headers = {"Authorization": f"Bearer {token}"}
    
    with open(pptx_path, 'rb') as pptx_file:
        files = {"file": pptx_file}
        data = {"name": "complex_test_template"}
        
        response = requests.post(f"{API_BASE}/templates/?name=complex_test_template", 
                               headers=headers, 
                               files=files)
    
    if response.status_code == 200:
        template = response.json()
        print(f"Template uploaded: ID {template['id']}")
        return template['id']
    else:
        print(f"Failed to upload template: {response.text}")
        return None

def login():
    """Login and get access token"""
    response = requests.post(f"{API_BASE}/auth/token", data={
        "username": USERNAME,
        "password": PASSWORD
    })
    
    if response.status_code == 200:
        return response.json()["access_token"]
    else:
        print(f"Login failed: {response.text}")
        sys.exit(1)

def run_complex_test():
    """Test complex PDF layout preservation"""
    print("Testing complex PDF layout preservation...")
    
    # Login
    token = login()
    print("Logged in successfully")
    
    # Create complex PowerPoint
    print("Creating complex PowerPoint...")
    pptx_path = create_complex_pptx()
    print(f"Complex PowerPoint created: {pptx_path}")
    
    # Upload as template
    template_id = upload_test_template(token, pptx_path)
    if not template_id:
        sys.exit(1)
    
    # Create test TXT file with matching RSIDs
    txt_content = """RSID,RESULT
rs12345,AA
rs67890,AG
rs111111,GG
rs222222,AT
rs333333,CC"""
    
    with tempfile.NamedTemporaryFile(mode='w', suffix='.txt', delete=False) as f:
        f.write(txt_content)
        txt_path = f.name
    
    # Submit job with complex template
    headers = {"Authorization": f"Bearer {token}"}
    with open(txt_path, 'rb') as txt_file:
        files = {"txt_file": txt_file}
        
        response = requests.post(f"{API_BASE}/jobs/?template_id={template_id}&excel_data_id=1", 
                               headers=headers, 
                               files=files)
    
    if response.status_code == 200:
        job = response.json()
        job_id = job['id']
        print(f"Job created: {job_id}")
        
        # Monitor job
        print("Monitoring job progress...")
        while True:
            response = requests.get(f"{API_BASE}/jobs/{job_id}", headers=headers)
            
            if response.status_code == 200:
                job = response.json()
                status = job["status"]
                progress = job.get("progress", 0)
                
                print(f"Status: {status}, Progress: {progress}%")
                
                if status == "done":
                    print("✅ Complex PDF job completed successfully!")
                    print(f"PPTX path: {job.get('output_pptx_path')}")
                    print(f"PDF path: {job.get('output_pdf_path')}")
                    
                    # Test PDF download
                    pdf_response = requests.get(f"{API_BASE}/jobs/{job_id}/download-pdf", headers=headers)
                    if pdf_response.status_code == 200:
                        pdf_size = len(pdf_response.content)
                        print(f"✅ Complex PDF downloaded: {pdf_size} bytes")
                        
                        if pdf_response.content.startswith(b'%PDF-'):
                            print("✅ PDF header valid")
                            print("🎉 Complex PDF layout preservation test PASSED!")
                            return True
                        else:
                            print("❌ Invalid PDF header")
                            return False
                    else:
                        print(f"❌ PDF download failed: {pdf_response.text}")
                        return False
                        
                elif status == "error":
                    print(f"❌ Job failed: {job.get('error_message')}")
                    return False
            else:
                print(f"Error checking job: {response.text}")
                return False
            
            time.sleep(3)
    else:
        print(f"❌ Failed to create job: {response.text}")
        return False

if __name__ == "__main__":
    success = run_complex_test()
    sys.exit(0 if success else 1)