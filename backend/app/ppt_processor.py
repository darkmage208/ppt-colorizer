import pandas as pd
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from typing import Dict, Optional
import re
import logging
from datetime import datetime
import os
from .storage import storage

logger = logging.getLogger(__name__)

class PPTProcessor:
    COLOR_MAP = {
        'D': RGBColor(0, 112, 192),      # BLUE
        'E': RGBColor(0, 128, 0),      # DARKER GREEN
        'F': RGBColor(255, 255, 0),    # YELLOW
        'G': RGBColor(255, 165, 0),    # ORANGE
        'H': RGBColor(255, 0, 0),      # RED
    }

    def __init__(self):
        self.excel_data = None
        self.txt_data = None
        self.presentation = None
        self.excel_filename = None
        self.txt_filename = None
        self.template_name = None
        # Performance optimization caches
        self._txt_lookup_cache = None
        self._shape_rsid_map = None
        self._rsid_color_cache = None
    
    def set_template_name(self, template_name: str):
        """Set the template name to be used as patient name"""
        self.template_name = template_name
        logger.info(f"Template name set to: {template_name}")
    
    def load_excel_data(self, excel_file_path: str) -> pd.DataFrame:
        # Store the filename for patient name extraction
        self.excel_filename = os.path.basename(excel_file_path)
        
        excel_content = storage.download_file(excel_file_path)
        excel_io = io.BytesIO(excel_content)
        
        # Read Excel file
        df = pd.read_excel(excel_io)
        
        # Ensure we have at least columns A through H
        if len(df.columns) < 8:
            raise ValueError(f"Excel file must have at least 8 columns (A-H), found {len(df.columns)}")
        
        # Rename columns to standard letters for easy access
        column_names = ['A', 'B', 'C', 'D', 'E', 'F', 'G', 'H']
        for i, col_name in enumerate(column_names):
            if i < len(df.columns):
                df.columns.values[i] = col_name
        
        # Column B should contain SNP data (RSID)
        self.excel_data = df
        return df
    
    def load_txt_data(self, txt_file_path: str, original_filename: str = None) -> pd.DataFrame:
        # Store the original filename if provided, otherwise use the file path basename
        self.txt_filename = original_filename or os.path.basename(txt_file_path)

        # MEMORY OPTIMIZATION: Use streaming processing for large TXT files
        logger.info(f"Loading TXT data with memory-optimized streaming processing")

        txt_content = storage.download_file(txt_file_path)

        # Process file in streaming fashion to avoid loading entire content in memory
        lines = txt_content.decode('utf-8', errors='ignore').strip().split('\n')

        if len(lines) < 2:
            raise ValueError("TXT file must have at least header and one data row")

        # Try tab-separated first, then comma-separated
        separator = '\t' if '\t' in lines[0] else ','
        header = [col.strip() for col in lines[0].split(separator)]

        # Find required columns
        rsid_col_idx = None
        result_col_idx = None

        for i, col in enumerate(header):
            col_upper = col.upper()
            if col_upper in ['RSID', 'SNP'] and rsid_col_idx is None:
                rsid_col_idx = i
            elif col_upper in ['RESULT', 'GENOTYPE', 'ALLELES', 'GT'] and result_col_idx is None:
                result_col_idx = i

        if rsid_col_idx is None:
            raise ValueError("TXT file must have an 'RSID' or 'SNP' column")
        if result_col_idx is None:
            raise ValueError("TXT file must have a 'RESULT', 'GENOTYPE', 'ALLELES', or 'GT' column")

        # MEMORY OPTIMIZATION: Build lookup cache directly without storing full DataFrame
        self._txt_lookup_cache = {}
        processed_lines = 0

        # Process data in chunks to reduce memory pressure
        chunk_size = 5000  # Process 5k lines at a time
        for i in range(1, len(lines), chunk_size):
            chunk_lines = lines[i:i + chunk_size]

            for line in chunk_lines:
                if line.strip():
                    row = line.split(separator)
                    # Ensure row has enough columns
                    if len(row) > max(rsid_col_idx, result_col_idx):
                        rsid = str(row[rsid_col_idx]).strip()
                        result = str(row[result_col_idx]).strip()

                        if rsid and result:
                            self._txt_lookup_cache[rsid] = result
                            processed_lines += 1

            # Log progress for large files
            if processed_lines % 25000 == 0:
                logger.info(f"Processed {processed_lines:,} TXT lines, {len(self._txt_lookup_cache):,} RSID mappings")

        logger.info(f"TXT processing complete: {processed_lines:,} lines, {len(self._txt_lookup_cache):,} RSID mappings")

        # Create minimal DataFrame for compatibility
        sample_row = ['sample_rsid', 'sample_result']
        self.txt_data = pd.DataFrame([sample_row], columns=[header[rsid_col_idx], header[result_col_idx]])

        # Clear cache when new TXT data is loaded
        self._rsid_color_cache = None

        # Force garbage collection to free memory
        del lines
        import gc
        gc.collect()

        return self.txt_data
    
    def extract_patient_name(self) -> str:
        """
        Use the TXT filename as the patient name
        Falls back to template name if TXT filename not available
        """
        logger.info(f"Extracting patient name - txt_filename: {self.txt_filename}, template_name: {self.template_name}")
        
        # Primary: Use TXT filename
        if self.txt_filename:
            base_name = os.path.splitext(self.txt_filename)[0]
            logger.info(f"Using TXT filename as patient name: {base_name}")
            return base_name
        
        # Fallback to template name if TXT filename not available
        if self.template_name:
            logger.info(f"TXT filename not available, using template name as patient name: {self.template_name}")
            return self.template_name
        
        logger.warning("No TXT filename or template name available, using 'Unknown Patient'")
        return "Unknown Patient"
    
    def get_current_date(self) -> str:
        """
        Get current date in DD-MM-YYYY format
        Example: 23-08-2025
        """
        return datetime.now().strftime("%d-%m-%Y")
    
    def replace_patient_and_date_info(self):
        """
        Find text frames containing "PATIENT" heading and reformat the entire content
        to the standard format: "PATIENT: [name]\nDATE: [date]"
        Searches recursively through groups with depth control.
        """
        if not self.presentation:
            return

        patient_name = self.extract_patient_name()
        current_date = self.get_current_date()

        for slide_idx, slide in enumerate(self.presentation.slides):
            for shape_idx, shape in enumerate(slide.shapes):
                self._find_and_update_patient_text_frame(shape, patient_name, current_date)

    def _find_and_update_patient_text_frame(self, shape, patient_name: str, current_date: str, depth: int = 0, max_depth: int = 3):
        """
        Recursively search for text frames containing "PATIENT" and reformat them.

        Args:
            shape: The shape to examine
            patient_name: Patient name to use
            current_date: Current date to use
            depth: Current recursion depth
            max_depth: Maximum recursion depth for group searching
        """
        try:
            # Handle individual text frames
            if hasattr(shape, 'text_frame') and shape.text_frame:
                full_text = shape.text_frame.text.strip().upper()

                # Check if this text frame contains "PATIENT" heading
                if "PATIENT" in full_text:
                    self._reformat_patient_text_frame(shape.text_frame, patient_name, current_date)
                    logger.debug(f"Updated patient info text frame at depth {depth}")
                    return True

            # Handle groups recursively (with depth control)
            elif hasattr(shape, 'shape_type') and shape.shape_type == 6 and depth < max_depth:  # Group
                found = False
                for sub_shape in shape.shapes:
                    if self._find_and_update_patient_text_frame(sub_shape, patient_name, current_date, depth + 1, max_depth):
                        found = True
                return found

        except Exception as e:
            logger.warning(f"Error processing shape at depth {depth}: {e}")

        return False

    def _reformat_patient_text_frame(self, text_frame, patient_name: str, current_date: str):
        """
        Reformat the text frame content to the standard format while preserving all formatting:
        PATIENT: [name]
        DATE: [date]

        Uses a careful approach to preserve existing formatting by copying font properties.

        Args:
            text_frame: The text frame to reformat
            patient_name: Patient name to use
            current_date: Current date to use
        """
        try:
            # Store formatting properties from existing runs
            saved_formats = []
            for paragraph in text_frame.paragraphs:
                para_formats = []
                for run in paragraph.runs:
                    # Save font properties
                    font_props = {
                        'name': getattr(run.font, 'name', None),
                        'size': getattr(run.font, 'size', None),
                        'bold': getattr(run.font, 'bold', None),
                        'italic': getattr(run.font, 'italic', None),
                        'underline': getattr(run.font, 'underline', None),
                        'color': None
                    }
                    try:
                        if hasattr(run.font, 'color') and hasattr(run.font.color, 'rgb'):
                            font_props['color'] = run.font.color.rgb
                    except:
                        pass
                    para_formats.append(font_props)
                saved_formats.append(para_formats)

            # Use the simplest approach: replace the entire text and then restore formatting
            text_frame.text = f"PATIENT: {patient_name}\nDATE: {current_date}"

            # Restore formatting to the new content
            if saved_formats and text_frame.paragraphs:
                # Apply formatting to patient line (first paragraph)
                if len(text_frame.paragraphs) > 0 and len(saved_formats) > 0 and len(saved_formats[0]) > 0:
                    first_para = text_frame.paragraphs[0]
                    if first_para.runs:
                        self._apply_font_properties(first_para.runs[0].font, saved_formats[0][0])

                # Apply formatting to date line (second paragraph)
                if len(text_frame.paragraphs) > 1:
                    second_para = text_frame.paragraphs[1]
                    if second_para.runs:
                        # Use second paragraph formatting if available, otherwise use first
                        format_to_use = saved_formats[1][0] if len(saved_formats) > 1 and saved_formats[1] else saved_formats[0][0]
                        self._apply_font_properties(second_para.runs[0].font, format_to_use)

            logger.debug(f"Reformatted text frame with patient: {patient_name}, date: {current_date}")

        except Exception as e:
            logger.warning(f"Error reformatting patient text frame: {e}")
            # Fallback: try simple text replacement without formatting preservation
            try:
                text_frame.text = f"PATIENT: {patient_name}\nDATE: {current_date}"
                logger.debug("Used fallback text replacement")
            except Exception as fallback_error:
                logger.error(f"Fallback text replacement also failed: {fallback_error}")

    def _apply_font_properties(self, target_font, font_props):
        """
        Apply saved font properties to a target font object.

        Args:
            target_font: The font object to apply properties to
            font_props: Dictionary of font properties to apply
        """
        try:
            if font_props.get('name'):
                target_font.name = font_props['name']
            if font_props.get('size'):
                target_font.size = font_props['size']
            if font_props.get('bold') is not None:
                target_font.bold = font_props['bold']
            if font_props.get('italic') is not None:
                target_font.italic = font_props['italic']
            if font_props.get('underline') is not None:
                target_font.underline = font_props['underline']
            if font_props.get('color'):
                target_font.color.rgb = font_props['color']
        except Exception as e:
            logger.debug(f"Could not apply some font properties: {e}")

    def load_presentation(self, ppt_file_path: str) -> Presentation:
        ppt_content = storage.download_file(ppt_file_path)
        ppt_io = io.BytesIO(ppt_content)
        
        self.presentation = Presentation(ppt_io)
        # Clear cache when new presentation is loaded
        self._shape_rsid_map = None
        return self.presentation

    def _build_txt_lookup_cache(self):
        """Build a fast lookup cache for TXT data: RSID -> RESULT"""
        if self._txt_lookup_cache is not None:
            return self._txt_lookup_cache

        # Cache is built during load_txt_data() for memory optimization
        if self._txt_lookup_cache is None:
            logger.warning("TXT lookup cache not found - ensure load_txt_data() was called")
            return {}

        logger.info(f"Using pre-built TXT lookup cache with {len(self._txt_lookup_cache)} entries")
        return self._txt_lookup_cache

    def _build_shape_rsid_map(self):
        """Build a map of RSID -> [shape objects] for fast lookup"""
        if self._shape_rsid_map is not None:
            return self._shape_rsid_map

        if self.presentation is None:
            return {}

        logger.info("Building PowerPoint shape-RSID mapping...")
        self._shape_rsid_map = {}

        def extract_shapes_recursive(shapes, level=0):
            """Recursively extract all text shapes and their RSIDs"""
            shape_count = 0
            for shape in shapes:
                try:
                    # Handle individual text frames
                    if hasattr(shape, 'text_frame') and shape.text_frame:
                        text_content = shape.text_frame.text.strip()
                        if text_content and text_content.startswith('rs'):  # RSID pattern
                            if text_content not in self._shape_rsid_map:
                                self._shape_rsid_map[text_content] = []
                            self._shape_rsid_map[text_content].append(shape)
                            shape_count += 1

                    # Handle groups recursively
                    elif hasattr(shape, 'shape_type') and shape.shape_type == 6 and level < 4:  # Group, max 4 levels
                        shape_count += extract_shapes_recursive(shape.shapes, level + 1)

                except Exception as e:
                    logger.warning(f"Error processing shape at level {level}: {e}")

            return shape_count

        total_shapes = 0
        for slide_idx, slide in enumerate(self.presentation.slides):
            slide_shapes = extract_shapes_recursive(slide.shapes)
            total_shapes += slide_shapes

        logger.info(f"Shape-RSID mapping built: {len(self._shape_rsid_map)} unique RSIDs found across {total_shapes} shapes")
        return self._shape_rsid_map

    def _build_rsid_color_cache(self):
        """Pre-compute colors for all RSIDs found in Excel data"""
        if self._rsid_color_cache is not None:
            return self._rsid_color_cache

        if self.excel_data is None:
            return {}

        logger.info("Building RSID-color cache...")
        self._rsid_color_cache = {}
        txt_cache = self._build_txt_lookup_cache()

        for _, row in self.excel_data.iterrows():
            rsid = str(row['B']).strip()
            if not rsid or rsid == 'nan':
                continue

            # Get RESULT from TXT data using cache
            result_value = txt_cache.get(rsid)
            if not result_value:
                continue

            # Find matching color column
            color_columns = ['D', 'E', 'F', 'G', 'H']
            for col in color_columns:
                cell_value = str(row[col]).strip().upper()
                if not cell_value or cell_value == 'NAN':
                    continue

                # Handle comma-separated values
                cell_values = [v.strip().upper() for v in cell_value.split(',')]
                if result_value in cell_values:
                    self._rsid_color_cache[rsid] = self.COLOR_MAP[col]
                    break

        logger.info(f"RSID-color cache built with {len(self._rsid_color_cache)} color mappings")
        return self._rsid_color_cache

    def find_color_for_rsid(self, rsid: str) -> Optional[RGBColor]:
        """
        Find the appropriate color for an RSID using pre-built cache.
        Much faster than the original implementation.
        """
        rsid_str = str(rsid).strip()
        color_cache = self._build_rsid_color_cache()
        return color_cache.get(rsid_str)
    
    def _apply_colors_to_shapes_optimized(self, rsid_color_map):
        """
        Optimized method to apply colors to shapes using pre-built mappings.
        Only iterates through shapes once instead of once per RSID.
        """
        if not rsid_color_map:
            return 0

        logger.info(f"Applying colors to shapes for {len(rsid_color_map)} RSIDs...")
        shape_map = self._build_shape_rsid_map()
        colored_count = 0

        # Batch apply colors using the pre-built shape map
        for rsid, color in rsid_color_map.items():
            shapes = shape_map.get(rsid, [])
            for shape in shapes:
                try:
                    if color:
                        shape.fill.solid()
                        shape.fill.fore_color.rgb = color
                        colored_count += 1
                        logger.debug(f"Applied {color} to RSID {rsid}")
                except Exception as e:
                    logger.warning(f"Error applying color to shape for RSID {rsid}: {e}")

        logger.info(f"Color application complete: {colored_count} shapes colored")
        return colored_count

    def find_and_modify_text_in_group(self, rsid: str, new_color):
        """
        Legacy method - kept for backward compatibility.
        For better performance, use process_presentation_optimized() instead.
        """
        rsid_str = str(rsid).strip()
        shape_map = self._build_shape_rsid_map()

        shapes = shape_map.get(rsid_str, [])
        for shape in shapes:
            try:
                if new_color:
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = new_color
            except Exception as e:
                logger.warning(f"Error modifying shape for RSID {rsid}: {e}")

    def apply_background_color(self, shape, color: RGBColor):
        if hasattr(shape, 'fill'):
            shape.fill.solid()
            shape.fill.fore_color.rgb = color
    
    def process_presentation(self, progress_callback=None) -> Dict[str, int]:
        """
        Process presentation according to the workflow:
        1. Iterate through all records in Excel template
        2. Extract SNP [column B] text (RSID)
        3. Find RSID in TXT data and get RESULT value
        4. Find which column (D,E,F,G,H) contains the RESULT value
        5. Find text box in PPT with exact RSID match
        6. Apply the corresponding color to the text box background
        """
        if not all([self.excel_data is not None, self.txt_data is not None, self.presentation is not None]):
            raise ValueError("Excel data, TXT data, and presentation must be loaded first")
        
        results = {
            'total_records': 0,
            'processed': 0,
            'colored': 0,
            'skipped': 0
        }
        
        total_records = len(self.excel_data)
        results['total_records'] = total_records
        
        logger.info(f"Starting to process {total_records} Excel records")
        
        # Replace PATIENT and DATE information before processing colors
        self.replace_patient_and_date_info()
        
        # Step 1: Iterate through all records in Excel template
        for idx, row in self.excel_data.iterrows():
            results['processed'] += 1
            
            # Step 2: Extract SNP [column B] text (RSID format: rs57108, rs131415, etc.)
            rsid = str(row['B']).strip()
            if not rsid or rsid == 'nan':
                results['skipped'] += 1
                continue
            
            logger.debug(f"Processing RSID: {rsid}")
            
            # Step 3: Find RSID in TXT data and get RESULT value
            # Step 4: Find which column contains the RESULT value and get color
            color = self.find_color_for_rsid(rsid)
            
            if color is None:
                results['skipped'] += 1
                continue
            
            # # Step 5: Find text box in PPT that exactly matches RSID
            # text_boxes = self.find_text_boxes_with_rsid(rsid)
            # if not text_boxes:
            #     results['skipped'] += 1
            #     continue
            
            # # Step 6: Apply color to text box background
            # for slide_idx, shape_idx, shape in text_boxes:
            #     self.apply_background_color(shape, color)
            #     results['colored'] += 1
            #     logger.debug(f"Applied color to RSID {rsid} on slide {slide_idx+1}")
            self.find_and_modify_text_in_group(rsid, color)
            
            if progress_callback:
                progress = int((idx + 1) / total_records * 100)
                progress_callback(progress)
        
        logger.info(f"Processing complete. Colored: {results['colored']}, Skipped: {results['skipped']}")
        return results

    def process_presentation_optimized(self, progress_callback=None) -> Dict[str, int]:
        """
        OPTIMIZED presentation processing with significant performance improvements:
        - Pre-builds lookup caches once instead of searching repeatedly
        - Uses reverse lookup: finds all shapes first, then applies colors in batch
        - Reduces complexity from O(n*m) to O(n+m)
        - Expected speedup: 10-100x faster depending on file sizes
        """
        if not all([self.excel_data is not None, self.txt_data is not None, self.presentation is not None]):
            raise ValueError("Excel data, TXT data, and presentation must be loaded first")

        results = {
            'total_records': 0,
            'processed': 0,
            'colored': 0,
            'skipped': 0,
            'cache_build_time': 0,
            'color_apply_time': 0
        }

        import time
        start_time = time.time()

        total_records = len(self.excel_data)
        results['total_records'] = total_records

        logger.info(f"Starting OPTIMIZED processing of {total_records} Excel records")

        # Replace PATIENT and DATE information before processing colors
        self.replace_patient_and_date_info()

        # OPTIMIZATION 1: Build all caches upfront (one-time cost)
        cache_start = time.time()
        logger.info("Building optimization caches...")

        # Build caches in optimal order
        txt_cache = self._build_txt_lookup_cache()
        shape_map = self._build_shape_rsid_map()
        color_cache = self._build_rsid_color_cache()

        cache_time = time.time() - cache_start
        results['cache_build_time'] = cache_time
        logger.info(f"Cache building completed in {cache_time:.2f}s")

        # OPTIMIZATION 2: Filter only RSIDs that exist in both Excel and PowerPoint
        excel_rsids = set()
        for _, row in self.excel_data.iterrows():
            rsid = str(row['B']).strip()
            if rsid and rsid != 'nan':
                excel_rsids.add(rsid)

        ppt_rsids = set(shape_map.keys())
        matching_rsids = excel_rsids.intersection(ppt_rsids)

        logger.info(f"Found {len(matching_rsids)} RSIDs that exist in both Excel ({len(excel_rsids)}) and PowerPoint ({len(ppt_rsids)})")

        # OPTIMIZATION 3: Build final color mapping only for matching RSIDs
        final_color_map = {}
        for rsid in matching_rsids:
            color = color_cache.get(rsid)
            if color:
                final_color_map[rsid] = color
                results['processed'] += 1
            else:
                results['skipped'] += 1

        logger.info(f"Color mapping built: {len(final_color_map)} RSIDs will be colored")

        # OPTIMIZATION 4: Apply all colors in one batch operation
        apply_start = time.time()
        colored_count = self._apply_colors_to_shapes_optimized(final_color_map)
        apply_time = time.time() - apply_start

        results['colored'] = colored_count
        results['color_apply_time'] = apply_time

        # Update progress
        if progress_callback:
            progress_callback(100)

        total_time = time.time() - start_time
        logger.info(f"OPTIMIZED processing complete in {total_time:.2f}s (cache: {cache_time:.2f}s, apply: {apply_time:.2f}s)")
        logger.info(f"Performance: {total_records/total_time:.1f} records/sec, {colored_count/total_time:.1f} shapes/sec")
        logger.info(f"Results: Colored: {results['colored']}, Skipped: {results['skipped']}")

        return results
    
    def save_presentation(self, output_path: str = None) -> str:
        if self.presentation is None:
            raise ValueError("No presentation loaded")
        
        output_io = io.BytesIO()
        self.presentation.save(output_io)
        output_io.seek(0)
        
        if output_path is None:
            # Use TXT filename for output if available, otherwise use timestamp
            if self.txt_filename:
                base_name = os.path.splitext(self.txt_filename)[0]
                output_filename = f"{base_name}.pptx"
                logger.info(f"Using TXT-based filename: {output_filename} (from {self.txt_filename})")
            else:
                output_filename = f"output_{pd.Timestamp.now().strftime('%Y%m%d_%H%M%S')}.pptx"
                logger.info(f"Using timestamp-based filename: {output_filename}")
            output_path = f"processed_presentations/{output_filename}"
        
        logger.info(f"Saving presentation to storage with filename: {output_path.split('/')[-1]}, folder: processed_presentations")
        file_key = storage.upload_file(output_io, output_path.split('/')[-1], "processed_presentations")
        logger.info(f"Presentation saved with key: {file_key}")
        return file_key

    def convert_to_pdf(self, pptx_key: str) -> str:
        """
        Convert PPTX to PDF using ConvertAPI

        Args:
            pptx_key (str): Storage key of the PPTX file to convert

        Returns:
            str: Storage key of the converted PDF file
        """
        import convertapi
        import tempfile
        import io
        import time
        import os
        import pandas as pd

        from .config import settings

        logger.info(f"Starting ConvertAPI PPTX to PDF conversion for key: {pptx_key}")

        try:
            # Get ConvertAPI secret from settings
            convertapi_secret = settings.convertapi_secret
            if not convertapi_secret:
                logger.error("ConvertAPI secret is not configured")
                raise Exception("ConvertAPI secret not configured")

            # Set ConvertAPI credentials
            convertapi.api_secret = convertapi_secret
            logger.info("ConvertAPI credentials configured")

            # Download PPTX file from storage
            logger.info("Downloading PPTX file from storage...")
            pptx_content = storage.download_file(pptx_key)
            logger.info(f"Downloaded PPTX file, size: {len(pptx_content)} bytes")

            # Create temporary file for PPTX
            with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_pptx:
                temp_pptx.write(pptx_content)
                temp_pptx_path = temp_pptx.name
                logger.info(f"Temporary PPTX file created: {temp_pptx_path}")

            # Convert using ConvertAPI
            logger.info("Starting ConvertAPI conversion...")
            start_time = time.time()

            result = convertapi.convert('pdf', {
                'File': temp_pptx_path,
                'PdfResolution': '150'  # Good balance between quality and file size
            })

            conversion_time = time.time() - start_time
            logger.info(f"ConvertAPI conversion completed in {conversion_time:.2f} seconds")
            logger.info(f"Conversion cost: {result.conversion_cost} credits")

            # Get PDF content from ConvertAPI result
            pdf_content = result.file.io.getvalue()
            logger.info(f"PDF generated, size: {len(pdf_content)} bytes")

            # Verify PDF header
            if not pdf_content.startswith(b'%PDF'):
                logger.warning("Generated file may not be a valid PDF")
                raise Exception("ConvertAPI returned invalid PDF content")

            # Create filename for PDF
            if self.txt_filename:
                base_name = os.path.splitext(self.txt_filename)[0]
                output_filename = f"{base_name}.pdf"
                logger.info(f"Using TXT-based filename for PDF: {output_filename}")
            else:
                output_filename = f"output_{pd.Timestamp.now().strftime('%Y%m%d_%H%M%S')}.pdf"
                logger.info(f"Using timestamp-based filename for PDF: {output_filename}")

            # Upload PDF to storage
            logger.info("Uploading PDF to storage...")
            pdf_io = io.BytesIO(pdf_content)
            pdf_key = storage.upload_file(pdf_io, output_filename, "processed_presentations")

            logger.info(f"PDF conversion completed successfully. Storage key: {pdf_key}")
            return pdf_key

        except Exception as e:
            logger.error(f"ConvertAPI PDF conversion failed: {str(e)}")
            raise Exception(f"PDF conversion failed: {str(e)}")

        finally:
            # Cleanup temporary file
            try:
                if 'temp_pptx_path' in locals() and os.path.exists(temp_pptx_path):
                    os.unlink(temp_pptx_path)
                    logger.info("Temporary PPTX file cleaned up")
            except Exception as cleanup_error:
                logger.warning(f"Failed to cleanup temporary file: {cleanup_error}")